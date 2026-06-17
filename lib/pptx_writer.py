#!/usr/bin/env python3
# ============================================================================
# designbook · lib/pptx_writer.py — deck IR -> editable .pptx (native shapes)
# ----------------------------------------------------------------------------
# IN-HOUSE authorship. Uses python-pptx (MIT). NO code, structure, or strings
# from GordenSuperPPTSkills (which has no license). Their pipeline extracts image
# LAYERS from a rendered PPT; ours is the opposite — the input is STRUCTURED deck
# HTML, so we map DOM -> native PowerPoint text boxes / shapes directly. Every
# slide is fully editable (real text runs, real shapes), never a flattened image.
#
# Honesty: this is STRUCTURALLY faithful, not pixel-identical. CSS gradients,
# backdrop-blur, web fonts and container-query sizing have no native .pptx
# equivalent. The `--verify` round-trip proves the artifact OPENS and is
# structurally complete (slide/shape/notes census) — never that it "looks the
# same". That census is the verification-moat guarantee.
#
#   build:  python pptx_writer.py <ir.json> <out.pptx>   -> prints census JSON
#   verify: python pptx_writer.py --verify <file.pptx>   -> prints census JSON
# ============================================================================
import sys, json
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IN = 914400
W = int(13.333 * IN)
H = int(7.5 * IN)            # 16:9
MX = int(0.9 * IN)          # content margin
CW = W - 2 * MX
HEAD_FONT = "Inter"          # falls back gracefully in PowerPoint


def rgb(hexstr, fallback="FFFFFF"):
    h = str(hexstr or "").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6 or any(c not in "0123456789abcdefABCDEF" for c in h):
        h = fallback
    return RGBColor.from_string(h.upper())


def text(slide, x, y, w, h, runs, size=18, bold=False, color="FFFFFF",
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, color, bold, size)]
    p = tf.paragraphs[0]
    p.alignment = align
    for (t, c, b, s) in runs:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.name = HEAD_FONT
        r.font.color.rgb = rgb(c)
    return tb


def bullets(slide, x, y, w, h, items, fg, accent, size=18):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        rd = p.add_run(); rd.text = "—  "
        rd.font.size = Pt(size); rd.font.bold = True; rd.font.name = HEAD_FONT; rd.font.color.rgb = rgb(accent)
        rt = p.add_run(); rt.text = it
        rt.font.size = Pt(size); rt.font.name = HEAD_FONT; rt.font.color.rgb = rgb(fg)
    return tb


def rect(slide, x, y, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def build(ir, out):
    prs = Presentation()
    prs.slide_width = Emu(W)
    prs.slide_height = Emu(H)
    th = ir.get("theme", {})
    BG = th.get("bg", "#0a0a12"); FG = th.get("fg", "#f2f3f7"); AC = th.get("accent", "#7c5cff")
    MU = th.get("muted", "#9aa0ac"); SF = th.get("surface", "#16161f"); ON = th.get("onAccent", "#0a0a12")
    blank = prs.slide_layouts[6]

    for sd in ir.get("slides", []):
        slide = prs.slides.add_slide(blank)
        rect(slide, 0, 0, W, H, BG)                      # full-bleed background
        layout = sd.get("layout", "bullets")

        if sd.get("kicker"):
            text(slide, MX, int(0.7 * IN), CW, int(0.4 * IN), sd["kicker"].upper(), size=13, bold=True, color=AC)

        if layout in ("title", "quote", "divider"):
            text(slide, MX, int(2.5 * IN), CW, int(2.2 * IN), sd.get("title", ""), size=44, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
            if sd.get("lead"):
                text(slide, MX, int(4.6 * IN), CW, int(1.2 * IN), sd["lead"], size=20, color=MU)
        elif layout == "cta":
            text(slide, MX, int(2.4 * IN), CW, int(1.4 * IN), sd.get("title", ""), size=40, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
            if sd.get("lead"):
                text(slide, MX, int(3.7 * IN), CW, int(0.8 * IN), sd["lead"], size=18, color=MU)
            rect(slide, MX, int(4.7 * IN), int(2.7 * IN), int(0.72 * IN), AC)
            text(slide, MX, int(4.7 * IN), int(2.7 * IN), int(0.72 * IN), sd.get("cta", "Get started"), size=16, bold=True, color=ON, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            text(slide, MX, int(0.95 * IN), CW, int(1.0 * IN), sd.get("title", ""), size=30, bold=True, color=FG)
            body_top = int(2.15 * IN)
            body_h = H - body_top - int(0.7 * IN)
            if layout == "stats":
                stats = sd.get("stats", [])
                if stats:
                    sw = CW // max(1, len(stats))
                    for i, st in enumerate(stats):
                        x = MX + i * sw
                        text(slide, x, body_top + int(0.5 * IN), sw, int(1.1 * IN), st.get("v", ""), size=52, bold=True, color=AC)
                        text(slide, x, body_top + int(1.7 * IN), sw, int(0.5 * IN), st.get("l", ""), size=15, color=MU)
            elif layout == "twocol":
                col_w = (CW - int(0.6 * IN)) // 2
                bullets(slide, MX, body_top, col_w, body_h, sd.get("bullets", []), FG, AC, size=18)
                rect(slide, MX + col_w + int(0.6 * IN), body_top, col_w, body_h, SF)
            elif layout == "media":
                if sd.get("lead"):
                    text(slide, MX, body_top, CW, int(0.6 * IN), sd["lead"], size=18, color=MU)
                    body_top += int(0.7 * IN); body_h -= int(0.7 * IN)
                rect(slide, MX, body_top, CW, body_h, SF)
            else:
                bullets(slide, MX, body_top, CW, body_h, sd.get("bullets", []), FG, AC, size=20)

        if sd.get("notes"):
            slide.notes_slide.notes_text_frame.text = sd["notes"]

    prs.save(out)


def census(path):
    prs = Presentation(path)
    slides = list(prs.slides)
    shapes = sum(len(s.shapes) for s in slides)
    pics = sum(1 for s in slides for sh in s.shapes if sh.shape_type == 13)
    notes = sum(1 for s in slides if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    return {"opens": True, "slides": len(slides), "shapes": shapes, "pictures": pics, "withNotes": notes}


if __name__ == "__main__":
    try:
        if len(sys.argv) >= 3 and sys.argv[1] == "--verify":
            print(json.dumps(census(sys.argv[2])))
        else:
            ir = json.load(open(sys.argv[1]))
            build(ir, sys.argv[2])
            print(json.dumps(census(sys.argv[2])))   # build-time round-trip self-check
    except Exception as e:
        print(json.dumps({"opens": False, "error": str(e)}), file=sys.stderr)
        sys.exit(1)
